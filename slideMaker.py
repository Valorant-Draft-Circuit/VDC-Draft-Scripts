from dataStructurePullerThingy import populateFranchises, tierPulls, franchises, picks
from pptx import Presentation

def makeSlides():
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    populateFranchises("VDCContracts-Teams.csv")
    tierPulls("Season1VDCDraft_Board-Bot_CommandOutput.csv", "Contender")
    tierPulls("Season1VDCDraft_Board-Bot_CommandOutput.csv", "Advanced")
    tierPulls("Season1VDCDraft_Board-Bot_CommandOutput.csv", "Master")
    tierPulls("Season1VDCDraft_Board-Bot_CommandOutput.csv", "Premier")
    print(picks.keys())
    print(picks["Contender"].keys())
    for tier in picks.keys():
        title = prs.slides.add_slide(slide_layout)
        title.shapes.title.text = tier + " draft starts now"
        for pick in picks[tier].keys():
                slide = prs.slides.add_slide(slide_layout)
                slide.placeholders[1].text = str(tier) + ": Round " + str(picks[tier][pick]['Round']) + ", Pick " + str(picks[tier][pick]['Pick']) + "\nFranchise: " + str(franchises[picks[tier][pick]["GM"]]['Franchise']) + "\nTeam: " + str(franchises[picks[tier][pick]["GM"]][tier])
                slide.shapes.title.text = picks[tier][pick]["Player Selected"]
    
    prs.save('vdcDraft.pptx')
    return None
makeSlides()